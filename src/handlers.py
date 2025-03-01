import PyPDF2
import docx
from pptx import Presentation

def read_pdf(file_path):
  print("Reading PDF file: {}".format(file_path))
  with open(file_path, 'rb') as file:
    reader = PyPDF2.PdfReader(file)
    return ' '.join([page.extract_text() for page in reader.pages])

def read_docx(file_path):
  print("Reading DOCX file: {}".format(file_path))
  doc = docx.Document(file_path)
  return ' '.join([paragraph.text for paragraph in doc.paragraphs])


def read_pptx(file_path):
  print("Reading PPTX file: {}".format(file_path))
  prs = Presentation(file_path)
  return ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])


def read_txt(file_path):
  print("Reading TXT file: {}".format(file_path))
  with open(file_path, 'r', encoding  = 'utf-8') as file:
    return file.read()
  
